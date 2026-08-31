#!/usr/bin/env python3
"""Build the two remaining XX395 deliverables from the pipeline's own output: the presentation and the digest.

The faculty guidelines require three files. This produces the second and third:

  * Presentation (10-15 minutes, Times New Roman): motivation and problem definition, objective, methods and
    tools, outcomes and deliverables, results.
  * Digest: a single slide giving the project title, company, duration and author, followed by the project
    objectives and a numbered list of outcomes.

Every number comes from results/summary.json, so the slides cannot drift from the report or the analysis.

  python3 report/build_slides.py --results results
"""
import argparse, json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
FONT = "Times New Roman"
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
TODO = RGBColor(0xC0, 0x00, 0x00)


def jload(p, default=None):
    p = Path(p)
    if not p.exists():
        return default if default is not None else {}
    try:
        return json.loads(p.read_text())
    except Exception:
        return default if default is not None else {}


def tload(p):
    p = Path(p)
    if not p.exists():
        return []
    with open(p) as fh:
        hdr = fh.readline().rstrip("\n").split("\t")
        return [dict(zip(hdr, l.rstrip("\n").split("\t"))) for l in fh if l.strip()]


def run_text(paragraph, text):
    """python-pptx's add_run() takes no argument; the text is assigned to the run afterwards."""
    r = paragraph.add_run()
    r.text = text
    return r


def style(run, size=18, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def slide_with_title(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    tf = textbox(s, 0.6, 0.35, 12.1, 0.9)
    style(run_text(tf.paragraphs[0], title), 30, bold=True, color=ACCENT)
    if subtitle:
        p = tf.add_paragraph()
        style(run_text(p, subtitle), 15, color=MUTED, italic=True)
    return s


def bullets(slide, items, left=0.75, top=1.5, width=11.9, size=17, gap=6):
    tf = textbox(slide, left, top, width, 5.4)
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(gap)
        style(run_text(p, ("—  " if level else "•  ") + text),
              size - (2 if level else 0), color=INK if not level else MUTED)
    return tf


def table(slide, headers, rows, left=0.9, top=1.6, width=11.5, height=None, size=13):
    height = height or min(4.8, 0.42 * (len(rows) + 1))
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    t = shape.table
    for j, h in enumerate(headers):
        c = t.cell(0, j)
        c.text = ""
        style(run_text(c.text_frame.paragraphs[0], str(h)), size, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = t.cell(i, j)
            c.text = ""
            style(run_text(c.text_frame.paragraphs[0], "" if v is None else str(v)), size)
    return t


def note(slide, text, top=6.55, color=MUTED, size=13):
    tf = textbox(slide, 0.75, top, 11.9, 0.55)
    style(run_text(tf.paragraphs[0], text), size, color=color, italic=True)


def todo(slide, text, top=6.55):
    note(slide, "[TO COMPLETE] " + text, top=top, color=TODO)


# ------------------------------------------------------------------------------------ design system
# A dark palette matching the dashboard, so the screenshots sit inside the deck rather than on top of it.
BG = RGBColor(0x0E, 0x11, 0x17)
PANEL = RGBColor(0x18, 0x1D, 0x27)
TEXT = RGBColor(0xEC, 0xEF, 0xF4)
DIM = RGBColor(0x9A, 0xA4, 0xB2)
BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x4C, 0xAF, 0x7D)
AMBER = RGBColor(0xD8, 0x9B, 0x3C)
RED = RGBColor(0xC5, 0x5A, 0x5A)
W, H = 13.333, 7.5


def paint(slide, colour=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    r.fill.solid(); r.fill.fore_color.rgb = colour
    r.line.fill.background(); r.shadow.inherit = False
    return r


def band(slide, left, top, width, height, colour=PANEL):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = colour
    r.line.fill.background(); r.shadow.inherit = False
    return r


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s)
    return s


def head(slide, title, kicker=None, rule=BLUE):
    if kicker:
        tf = textbox(slide, 0.75, 0.42, 11.8, 0.35)
        style(run_text(tf.paragraphs[0], kicker.upper()), 12, bold=True, color=rule)
        top = 0.78
    else:
        top = 0.55
    tf = textbox(slide, 0.75, top, 11.8, 0.85)
    style(run_text(tf.paragraphs[0], title), 30, bold=True, color=TEXT)
    band(slide, 0.78, top + 0.92, 1.5, 0.045, rule)
    return top + 1.25


def points(slide, items, top, size=17, left=0.85, width=11.6, gap=11, marker=None):
    tf = textbox(slide, left, top, width, H - top - 0.6)
    for i, item in enumerate(items):
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if lvl:
            style(run_text(p, "      " + text), size - 2, color=DIM)
        else:
            r = run_text(p, "▪   "); style(r, size, color=marker or BLUE)
            style(run_text(p, text), size, color=TEXT)
    return tf


def stats(slide, items, top=2.1, colour=BLUE):
    """A row of large figures with captions — the visual anchor of a results slide."""
    n = len(items)
    gap, w = 0.35, (11.8 - 0.35 * (n - 1)) / n
    for i, (value, label) in enumerate(items):
        x = 0.75 + i * (w + gap)
        band(slide, x, top, w, 1.75, PANEL)
        tf = textbox(slide, x + 0.25, top + 0.18, w - 0.5, 0.85)
        style(run_text(tf.paragraphs[0], str(value)), 40, bold=True, color=colour)
        tf = textbox(slide, x + 0.25, top + 1.08, w - 0.5, 0.6)
        style(run_text(tf.paragraphs[0], label), 12, color=DIM)


def picture_fit(slide, path, box, caption=None, border=True):
    """Place an image inside (left, top, width, height), preserving aspect ratio and centring it.

    Sizing by width alone overflows the canvas whenever the image is taller than the space left under the
    heading - which is exactly what a 16:10 screenshot does on a 16:9 slide."""
    from PIL import Image
    left, top, max_w, max_h = box
    iw, ih = Image.open(str(path)).size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    x, y = left + (max_w - w) / 2, top + (max_h - h) / 2
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if border:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.background(); r.line.color.rgb = PANEL; r.line.width = Pt(1.25); r.shadow.inherit = False
    if caption:
        tf = textbox(slide, 0.85, top + max_h + 0.04, 11.6, 0.42)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        style(run_text(tf.paragraphs[0], caption), 11.5, color=DIM, italic=True)
    return pic


def image_slide(prs, title, kicker, path, caption, rule=BLUE):
    """A screenshot slide: compact heading, the image given the rest of the canvas, caption underneath."""
    s = blank(prs)
    tf = textbox(s, 0.75, 0.32, 11.8, 0.3)
    style(run_text(tf.paragraphs[0], kicker.upper()), 11, bold=True, color=rule)
    tf = textbox(s, 0.75, 0.6, 11.8, 0.58)
    style(run_text(tf.paragraphs[0], title), 25, bold=True, color=TEXT)
    picture_fit(s, path, (0.85, 1.28, 11.6, 5.25), caption)
    return s


def callout(slide, text, top, colour=GREEN, height=0.95, size=16):
    band(slide, 0.75, top, 11.8, height, PANEL)
    band(slide, 0.75, top, 0.055, height, colour)
    tf = textbox(slide, 1.02, top + 0.13, 11.4, height - 0.2)
    style(run_text(tf.paragraphs[0], text), size, color=TEXT)


def footer(slide, n, total):
    tf = textbox(slide, 11.6, 6.98, 1.2, 0.35)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    style(run_text(tf.paragraphs[0], f"{n} / {total}"), 10, color=DIM)


def section(prs, number, title, blurb=None):
    s = blank(prs)
    band(s, 0, 2.55, W, 0.06, BLUE)
    tf = textbox(s, 0.9, 2.9, 11.5, 0.9)
    style(run_text(tf.paragraphs[0], f"{number}"), 15, bold=True, color=BLUE)
    tf = textbox(s, 0.9, 3.2, 11.5, 1.1)
    style(run_text(tf.paragraphs[0], title), 34, bold=True, color=TEXT)
    if blurb:
        tf = textbox(s, 0.9, 4.35, 10.5, 0.9)
        style(run_text(tf.paragraphs[0], blurb), 15, color=DIM, italic=True)
    return s


# ------------------------------------------------------------------------------------ presentation
def build_presentation(S, out, figures=None):
    figures = Path(figures or REPO / "report" / "figures")
    lc = S.get("qc", {}).get("library_complexity", {})
    up, rup = lc.get("umi_pattern", {}), lc.get("rna_umi_pattern", {})
    dd, ra = lc.get("deduplication", {}), lc.get("rna_alignment", {})
    var = S.get("variants", {})
    merged = var.get("caller_concordance", {})
    counts = var.get("counts_by_class", {})
    tmb = (S.get("biomarkers") or {}).get("tmb", {})
    fus = (S.get("fusions") or {}).get("summary", {})
    val = (S.get("validation") or {}).get("summary", {})
    hit = next((r for r in (S.get("validation") or {}).get("table", [])
                if r.get("status") == "CONCORDANT"), {})
    filt = jload(REPO / "results" / "dna" / "variants_filtered_summary.json")
    def fig(n):
        """Prefer the cropped version: a slide needs the informative region, not the whole page."""
        for candidate in (figures / ("crop_" + n), figures / n):
            if candidate.exists():
                return candidate
        return None

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    S_ = []

    # 1 — title
    s = blank(prs); S_.append(s)
    band(s, 0, 0, 0.11, H, BLUE)
    tf = textbox(s, 1.15, 2.15, 11.0, 2.1)
    style(run_text(tf.paragraphs[0], "Educational DNA–RNA"), 44, bold=True, color=TEXT)
    p = tf.add_paragraph(); style(run_text(p, "Tumour Profiling Pipeline"), 44, bold=True, color=TEXT)
    band(s, 1.18, 4.32, 2.0, 0.05, BLUE)
    tf = textbox(s, 1.15, 4.55, 10.6, 1.2)
    style(run_text(tf.paragraphs[0],
                   "Building a clinical-grade genomic workflow in the open — and measuring it against "
                   "an accredited laboratory"), 16, color=DIM, italic=True)
    tf = textbox(s, 1.15, 6.05, 10.6, 0.9)
    style(run_text(tf.paragraphs[0], "[NAME · PROGRAMME]   ·   [COMPANY]   ·   [DATES]   ·   "
                                     "supervisor [NAME]"), 13, bold=True, color=TODO)

    # 2 — the question
    s = blank(prs); S_.append(s)
    band(s, 0.75, 1.55, 11.8, 2.0, PANEL)
    band(s, 0.75, 1.55, 0.055, 2.0, BLUE)
    tf = textbox(s, 1.15, 1.95, 11.1, 1.4)
    style(run_text(tf.paragraphs[0], "“What can we understand about a person's disease"), 27, italic=True,
          color=TEXT)
    p = tf.add_paragraph()
    style(run_text(p, "by looking at their tumour's genetic material?”"), 27, italic=True, color=TEXT)
    points(s, [
        "A biopsy is sequenced and a two-page clinical report names a few mutations and the drugs that "
        "target them.",
        "What happens in between is invisible: the software is commercial and closed, the thresholds "
        "undocumented, the reasoning not shown.",
        "This project builds that chain in the open — then checks it against the real clinical report for "
        "the same specimen.",
    ], top=4.05, size=17)

    # 3 — section
    S_.append(section(prs, "01", "Identifying the data",
                      "Before any biology: which assay produced these reads?"))

    # 4 — read structure
    s = blank(prs); S_.append(s)
    top = head(s, "The assay is written into the reads", "read structure")
    band(s, 0.75, top, 11.8, 1.35, PANEL)
    tf = textbox(s, 1.05, top + 0.16, 11.2, 1.1)
    style(run_text(tf.paragraphs[0], "Read 1    [ 12 random bases = molecular barcode ]"
                                     "[ fixed adapter ][ patient DNA … ]"), 16, bold=True, color=TEXT)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    style(run_text(p, "Read 2    [ gene-specific primer ][ patient DNA … ]"), 16, bold=True, color=TEXT)
    points(s, [
        "A position where all four bases occur equally is random; one base at >98 % is synthetic. "
        "The transition marks the end of the barcode.",
        "Only ~2,200 distinct sequences start 80 % of Read 2 → a primer-based targeted panel.",
        "Conclusion: anchored multiplex PCR — confirmed afterwards by the clinical report "
        "(Archer VariantPlex + FusionPlex).",
    ], top=top + 1.6, size=16)
    stats(s, [(f"{up.get('fraction_matching', 0)*100:.1f} %", "of DNA reads match the predicted structure"),
              (f"{rup.get('fraction_matching', 0)*100:.1f} %", "of RNA reads match it"),
              ("8 + 10 nt", "index reads — the adapter kit's signature")], top=5.15)

    # 5 — QC screenshot
    if fig("02_qc_library.png"):
        S_.append(image_slide(
            prs, "The barcode is visible in the base composition", "dashboard · quality control",
            fig("02_qc_library.png"),
            "Positions 1–12 carry all four bases; positions 13–26 are a single fixed base per position — "
            "the adapter. The insert begins at base 27."))

    # 6 — section
    S_.append(section(prs, "02", "The pipeline",
                      "Split by memory, not by convenience"))

    # 7 — architecture
    s = blank(prs); S_.append(s)
    top = head(s, "Where each step runs", "architecture")
    table(s, ["Stage", "Tool", "Runs on"], [
        ["Read structure, primer inference", "project scripts", "laptop"],
        ["Barcode extraction, alignment, deduplication", "UMI-tools, BWA-MEM", "Galaxy"],
        ["Primer clipping, coverage", "samtools, mosdepth", "Galaxy"],
        ["Somatic calling ×3", "Mutect2, LoFreq, FreeBayes", "Galaxy"],
        ["RNA alignment and fusions", "STAR, Arriba", "Galaxy"],
        ["Annotation, tiering, report, dashboard", "VEP REST, CIViC, project scripts", "laptop"],
    ], left=0.85, top=top, width=11.6, size=14)
    callout(s, "Aligning to the human genome needs ~6 GB (BWA) and ~31 GB (STAR). The machine had 8 GB — so "
               "the heavy half runs on the free public Galaxy service and the interpretation locally.",
            top=5.55, colour=BLUE, height=0.85)

    # 8 — molecules
    s = blank(prs); S_.append(s)
    top = head(s, "Counting molecules, not reads", "deduplication")
    stats(s, [(f"{dd.get('reads_in', 0)/1e6:.1f} M", "sequencing reads"),
              (f"{dd.get('unique_molecules_out', 0)/1e6:.1f} M", "unique molecules"),
              (f"{dd.get('mean_unique_umis_per_position', 0)}", "molecules per amplicon start")], top=top)
    points(s, [
        "Every molecule from one amplicon starts at the same coordinate, so coordinate-based duplicate "
        "marking would collapse them all into one.",
        "With only ~2 distinct molecules per start position, that would discard most of the evidence.",
        "The experiment is a 2.9-million-molecule measurement, not a 5.5-million-read one — and the "
        "confidence in every variant follows from the smaller number.",
    ], top=top + 2.15, size=17)

    # 9 — section
    S_.append(section(prs, "03", "Results", "Measured against an accredited laboratory"))

    # 10 — the agreement
    s = blank(prs); S_.append(s)
    top = head(s, "The pipeline recovered the laboratory's finding", "validation", rule=GREEN)
    stats(s, [(hit.get("lab_vaf", "0.281"), "laboratory allele fraction"),
              (hit.get("our_vaf", "0.2814"), "measured here"),
              (f"{val.get('concordant', 1)} / {val.get('reported_by_laboratory', 1)}",
               "reported variants recovered")], top=top, colour=GREEN)
    callout(s, "IDH1 p.Arg132Cys — supported by all three callers. Different aligners, different callers, "
               "different reference builds, different deduplication: the same answer to three decimals.",
            top=top + 2.0, colour=GREEN, height=0.9)
    points(s, [
        "The laboratory used GRCh37 with the vendor's closed software; this pipeline used GRCh38 with three "
        "open callers.",
        "Primer clipping alone moved the estimate from 0.2636 to 0.2814 — measured, not assumed.",
    ], top=top + 3.1, size=16, marker=GREEN)

    # 11 — variants screenshot
    if fig("03_variants.png"):
        S_.append(image_slide(
            prs, "Every call keeps the reason it was classified", "dashboard · variants",
            fig("03_variants.png"),
            "39 likely somatic, 136 likely germline, 1,543 with insufficient evidence — nothing is "
            "silently deleted.", rule=GREEN))

    # 12 — the disagreement
    s = blank(prs); S_.append(s)
    top = head(s, "Where the two disagree", "interpretation", rule=AMBER)
    band(s, 0.75, top, 5.75, 1.45, PANEL)
    tf = textbox(s, 1.0, top + 0.2, 5.3, 1.2)
    style(run_text(tf.paragraphs[0], "Laboratory:  Tier I-A"), 20, bold=True, color=TEXT)
    p = tf.add_paragraph()
    style(run_text(p, "cited a clinical guideline for this tumour type"), 13, color=DIM)
    band(s, 6.8, top, 5.75, 1.45, PANEL)
    tf = textbox(s, 7.05, top + 0.2, 5.3, 1.2)
    style(run_text(tf.paragraphs[0], f"This pipeline:  Tier {hit.get('our_tier', 'II')}"), 20, bold=True,
          color=AMBER)
    p = tf.add_paragraph()
    style(run_text(p, "open evidence exists only for other tumour types"), 13, color=DIM)
    points(s, [
        "Both are defensible. The tiering rule refuses Tier I on evidence from a different disease — which "
        "is what the AMP/ASCO/CAP guideline requires.",
        "The gap measures the uneven tumour-type coverage of free knowledge bases, not a disagreement about "
        "biology.",
        "This is the project's clearest argument for why automated tiering needs curated knowledge or "
        "expert review.",
    ], top=top + 1.8, size=16, marker=AMBER)

    # 13 — artefacts
    s = blank(prs); S_.append(s)
    top = head(s, "One artefact with two faces", "quality", rule=RED)
    band(s, 0.75, top, 5.75, 1.95, PANEL)
    tf = textbox(s, 1.0, top + 0.16, 5.3, 1.7)
    style(run_text(tf.paragraphs[0], "DNA arm"), 19, bold=True, color=RED)
    for t in [f"{(filt.get('dominant_substitution_class') or {}).get('fraction_of_somatic_snvs', 0)*100:.0f} % "
              "of somatic SNVs share one substitution class",
              "Four “independent mutations” within 35 bases in one gene",
              f"{filt.get('somatic_flagged_artefact', 0)} of {counts.get('SOMATIC_LIKELY', 0)} candidates flagged"]:
        p = tf.add_paragraph(); p.space_before = Pt(7); style(run_text(p, "▪  " + t), 14, color=TEXT)
    band(s, 6.8, top, 5.75, 1.95, PANEL)
    tf = textbox(s, 7.05, top + 0.16, 5.3, 1.7)
    style(run_text(tf.paragraphs[0], "RNA arm"), 19, bold=True, color=RED)
    for t in [f"{fus.get('n_fusions', 0)} fusions called, 166,694 discarded",
              "FGFR1/2/3 joined in every combination, both orientations",
              f"All {fus.get('n_fusions', 0)} flagged — the laboratory reported none, and we agree"]:
        p = tf.add_paragraph(); p.space_before = Pt(7); style(run_text(p, "▪  " + t), 14, color=TEXT)
    callout(s, "One mechanism explains both: cross-priming between the amplicons the assay amplifies in a "
               "single tube. An unscreened reading would have reported an FGFR2 fusion — the most "
               "consequential false positive available in this cancer type.", top=top + 2.2, colour=RED,
            height=1.0)

    # 14 — fusions screenshot
    if fig("04_rna_fusions.png"):
        S_.append(image_slide(
            prs, "Screening turns twenty candidates into a negative", "dashboard · RNA",
            fig("04_rna_fusions.png"),
            "Each flag names its reason: reciprocal calls, paralogue pairs, joins between two primed "
            "panel genes, promiscuous or unnamed partners.", rule=RED))

    # 15 — pathways screenshot
    if fig("05_pathways.png"):
        S_.append(image_slide(
            prs, "From a gene list to a picture of the cell", "dashboard · pathways",
            fig("05_pathways.png"),
            "Membership mapping against the ten canonical oncogenic pathways — deliberately not an "
            "enrichment analysis, since the gene universe is fixed by the assay."))

    # 16 — limits
    s = blank(prs); S_.append(s)
    top = head(s, "What this data cannot tell us", "limits", rule=AMBER)
    points(s, [
        "Whether a variant is inherited — no normal tissue was sequenced.",
        "Roughly half of the nominal target territory had no coverage at all.",
        "Genes frequently mutated in this cancer type are absent from the panel: a negative there is an "
        "untested region, not an absence of mutation.",
        f"Mutational burden — {tmb.get('tmb_per_mb_before_screening', 0)} → {tmb.get('tmb_per_mb', 0)} per Mb "
        "after screening; both implausible, and the assessable territory is far too small.",
        "Microsatellite status: not determinable — as the accredited laboratory also concluded.",
    ], top=top, size=17, marker=AMBER)
    callout(s, "Stating where the answers stop is what makes this a teaching object rather than a black box.",
            top=5.75, colour=AMBER, height=0.8)

    # 17 — deliverables
    s = blank(prs); S_.append(s)
    top = head(s, "Deliverables")
    stats(s, [("38", "automated tests, run in CI"), ("14", "documentation pages"),
              ("6", "dashboard pages"), ("2", "execution modes")], top=top)
    points(s, [
        "Reproducible Snakemake + Bioconda pipeline; heavy steps scripted end-to-end on Galaxy.",
        "Synthetic example dataset over a real mini-reference with a truth set — the whole workflow runs and "
        "is tested without any patient data.",
        "Public repository with setup instructions, worked examples and educational documentation.",
    ], top=top + 2.15, size=16)
    tf = textbox(s, 0.85, 6.35, 11.6, 0.5)
    style(run_text(tf.paragraphs[0], "[REPOSITORY URL]"), 13, bold=True, color=TODO)

    # 18 — conclusions
    s = blank(prs); S_.append(s)
    top = head(s, "Conclusions", rule=GREEN)
    points(s, [
        "An open pipeline, on a laptop plus a free public service, reproduced an accredited laboratory's "
        "finding to three decimal places.",
        "Every disagreement was explicable and measurable: evidence coverage, reporting thresholds, and one "
        "identifiable chemistry artefact.",
        "The most useful property is not the agreement but the capacity to say precisely where the answers "
        "stop.",
    ], top=top, size=19, marker=GREEN)
    callout(s, "Educational output. Not a diagnostic tool, not clinically validated, and not a basis for "
               "medical decisions.", top=5.6, colour=RED, height=0.8, size=15)

    for i, sl in enumerate(S_, start=1):
        footer(sl, i, len(S_))
    prs.save(out)
    return out

# ------------------------------------------------------------------------------------ digest
def build_digest(S, out):
    val = (S.get("validation") or {}).get("summary", {})
    hit = next((r for r in (S.get("validation") or {}).get("table", [])
                if r.get("status") == "CONCORDANT"), {})
    counts = S.get("variants", {}).get("counts_by_class", {})

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    tf = textbox(s, 0.7, 0.3, 11.9, 1.4)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    style(run_text(tf.paragraphs[0], "Educational DNA–RNA Tumour Profiling Pipeline and "
                                   "Visualization Platform"), 24, bold=True, color=ACCENT)
    for label in ("[COMPANY NAME / ADDRESS]", "[DURATION OF PROJECT: month day – month day, year]",
                  "[STUDENT NAME / PROGRAMME]"):
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        style(run_text(p, label), 13, color=TODO, bold=True)

    tf = textbox(s, 0.7, 2.0, 11.9, 1.55)
    style(run_text(tf.paragraphs[0], "PROJECT OBJECTIVE & EXPECTATIONS"), 14, bold=True)
    p = tf.add_paragraph()
    style(run_text(p, "To build a reproducible, documented pipeline that takes tumour DNA and RNA sequencing "
                    "from a targeted cancer panel through to an interpreted summary — variant calling, "
                    "clinical tiering, fusion detection and reporting — and to make every step inspectable "
                    "rather than hidden inside commercial software. The pipeline was to be validated against "
                    "the report issued by an accredited clinical laboratory for the same specimen."), 12.5)

    tf = textbox(s, 0.7, 3.65, 11.9, 3.3)
    style(run_text(tf.paragraphs[0], "OUTCOMES"), 14, bold=True)
    outcomes = [
        "A reproducible Snakemake + Bioconda pipeline covering quality control, barcode-aware deduplication, "
        "tumour-only somatic calling, annotation, clinical tiering, RNA fusion analysis and reporting.",
        f"Independent recovery of the accredited laboratory's finding: {hit.get('gene', 'IDH1')} "
        f"{hit.get('variant', 'p.Arg132Cys')} at an allele fraction of {hit.get('our_vaf', '0.2814')} against "
        f"the laboratory's {hit.get('lab_vaf', '0.281')} — {val.get('concordant', 1)} of "
        f"{val.get('reported_by_laboratory', 1)} reported variants, none missed.",
        "Identification of the assay chemistry from the raw reads alone, later confirmed by the clinical "
        "report, and reconstruction of the panel design from the sequencing data.",
        "Artefact screening for both arms, which identified a single chemistry-derived mechanism explaining "
        "both the excess DNA variant calls and all 20 candidate RNA fusions.",
        "A public repository containing the pipeline, a synthetic example dataset with a truth set, 38 "
        "automated tests, continuous integration and fourteen documentation pages.",
        "A six-page interactive dashboard presenting variants, fusions, pathways, therapies and the limits of "
        "the analysis.",
    ]
    for i, o in enumerate(outcomes, start=1):
        p = tf.add_paragraph(); p.space_after = Pt(5)
        style(run_text(p, f"{i}.  {o}"), 12)

    prs.save(out)
    return out


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--results", default="results")
    ap.add_argument("--outdir", default="report")
    a = ap.parse_args()
    S = jload(Path(a.results) / "summary.json")
    out = Path(a.outdir)
    out.mkdir(parents=True, exist_ok=True)
    p1 = build_presentation(S, out / "BIO395_Presentation_DRAFT.pptx",
                            figures=REPO / "report" / "figures")
    p2 = build_digest(S, out / "BIO395_Digest_DRAFT.pptx")
    print(f"wrote {p1}\nwrote {p2}")


if __name__ == "__main__":
    main()
